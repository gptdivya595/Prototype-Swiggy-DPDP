#!/usr/bin/env python3
"""Build the light-mode Swiggy x DPDP Product Thinking Deck from docs/slides/deck.json."""
import json, os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DECK = json.load(open(os.path.join(ROOT, 'docs/slides/deck.json')))
DZ = os.path.join(ROOT, 'docs/diagrams')
DS = DECK['design_system']

def C(hex_):
    return RGBColor.from_string(hex_)

INK, MUTED, HAIR = C(DS['ink']), C(DS['muted']), C(DS['hairline'])
PRIM, PRIMD = C(DS['primary']), C(DS['primary_dark'])
GREEN, GREEN_BG = C(DS['green']), C(DS['green_bg'])
BLUE, BLUE_BG = C(DS['blue']), C(DS['blue_bg'])
RED, RED_BG = C(DS['red']), C(DS['red_bg'])
SURFACE, BG = C(DS['surface']), C(DS['bg'])
FH, FB = DS['font_head'], DS['font_body']

EMU_IN = 914400
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5
MX = 0.55  # side margin

# ---------- primitives ----------
def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def _set_font(run, size, color, bold=False, italic=False, font=FB):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        line=1.05, sp_after=2, wrap=True):
    """runs: list of paragraphs; each paragraph = list of (text,size,color,bold,italic,font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('left','right','top','bottom'):
        setattr(tf, 'margin_'+m, 0)
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.line_spacing = line; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        for seg in para:
            t, sz, col, b, it = seg[0], seg[1], seg[2], seg[3], seg[4]
            fnt = seg[5] if len(seg) > 5 else FB
            r = p.add_run(); r.text = t; _set_font(r, sz, col, b, it, fnt)
    return tb

def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, rounded=False, radius=0.06):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    return shp

def hline(s, x, y, w, color=HAIR, weight=1.0):
    ln = s.shapes.add_connector(2, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False
    return ln

def header(s, sl):
    rect(s, 0, 0, SW, 0.16, fill=PRIM)  # top accent band
    txt(s, MX, 0.42, SW-2*MX, 0.3, [[(sl['kicker'], 11, PRIMD, True, False)]])
    txt(s, MX, 0.72, SW-2*MX, 0.9, [[(sl['title'], 25, INK, True, False, FH)]], line=1.02)
    ty = 1.62
    if sl.get('subtitle'):
        txt(s, MX, ty, SW-2*MX, 0.6, [[(sl['subtitle'], 12.5, MUTED, False, False)]], line=1.08)
        ty += 0.5
    hline(s, MX, ty+0.02, SW-2*MX, HAIR, 1.0)
    return ty + 0.16

def footer(s, sl, n):
    hline(s, MX, SH-0.42, SW-2*MX, HAIR, 1.0)
    txt(s, MX, SH-0.36, 7, 0.3, [[(sl.get('footer',''), 9, MUTED, False, False)]])
    txt(s, SW-MX-4, SH-0.36, 4, 0.3,
        [[('Swiggy × DPDP · ', 9, MUTED, False, False), ('%02d / 10' % n, 9, PRIMD, True, False)]],
        align=PP_ALIGN.RIGHT)
    rect(s, SW-MX-0.02-0.0, SH-0.36, 0.0, 0.0)  # noop

def notes(s, sl):
    if sl.get('notes'):
        s.notes_slide.notes_text_frame.text = sl['notes']

def pic_fit(s, path, bx, by, bw, bh, align='center', valign='middle'):
    im = Image.open(path); iw, ih = im.size
    ar = iw/ih; bar = bw/bh
    if ar > bar:
        w = bw; h = bw/ar
    else:
        h = bh; w = bh*ar
    x = bx + (bw-w)/2 if align=='center' else bx
    y = by + (bh-h)/2 if valign=='middle' else by
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return x, y, w, h

def sev_color(sev):
    return {'Critical': (RED, RED_BG), 'High': (PRIMD, C('FCE9DE')),
            'Medium': (C('B45309'), C('FEF3C7'))}.get(sev, (MUTED, SURFACE))

def status_color(st):
    st = st.lower()
    if 'health' in st: return GREEN, GREEN_BG
    if 'crit' in st or 'below' in st: return RED, RED_BG
    return PRIMD, C('FCE9DE')

def card(s, x, y, w, h, heading, body, accent=PRIM, hcol=None, body_size=10.5, head_size=12):
    rect(s, x, y, w, h, fill=SURFACE, line=HAIR, lw=1.0, rounded=True, radius=0.05)
    rect(s, x, y+0.12, 0.06, h-0.24, fill=accent)
    txt(s, x+0.22, y+0.14, w-0.36, 0.4, [[(heading, head_size, hcol or INK, True, False, FH)]], line=1.0)
    if body:
        txt(s, x+0.22, y+0.14+0.36, w-0.4, h-0.5, [[(body, body_size, MUTED, False, False)]], line=1.08)

# ---------- bullet helper ----------
def bullets(s, x, y, w, h, items, size=11, color=INK, gap=3, marker='—  ', mcol=None):
    paras = [[(marker, size, mcol or PRIM, True, False), (it, size, color, False, False)] for it in items]
    return txt(s, x, y, w, h, paras, line=1.1, sp_after=gap)

# ---------- table helper ----------
def table(s, x, y, w, rows, colw, header_row=True, fs=10, hfs=10, rh=0.34, hh=0.36,
          cell_cols=None):
    ncol = len(colw); nrow = len(rows)
    tw = sum(colw)
    scale = w/tw; colw = [c*scale for c in colw]
    from pptx.util import Inches as I
    gtbl = s.shapes.add_table(nrow, ncol, I(x), I(y), I(w), I(hh + rh*(nrow-1))).table
    gtbl.first_row = False; gtbl.horz_banding = False
    for ci, cw in enumerate(colw):
        gtbl.columns[ci].width = I(cw)
    gtbl.rows[0].height = I(hh)
    for ri in range(1, nrow):
        gtbl.rows[ri].height = I(rh)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtbl.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_h = header_row and ri == 0
            if is_h:
                cell.fill.solid(); cell.fill.fore_color.rgb = PRIM
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = BG if ri % 2 else SURFACE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.line_spacing = 1.0
            r = p.add_run(); r.text = str(val)
            col = C('FFFFFF') if is_h else INK
            if cell_cols and not is_h and ci in cell_cols:
                col = cell_cols[ci](val)
            _set_font(r, hfs if is_h else fs, col, is_h, False)
    return gtbl

# ================= per-slide layouts =================
def s1(s, sl):
    y = header(s, sl)
    lw = 7.0
    # problem / solution stack
    yy = y + 0.1
    txt(s, MX, yy, lw, 0.3, [[('THE PROBLEM', 10, RED, True, False)]]); yy += 0.3
    txt(s, MX, yy, lw, 1.2, [[(sl['problem'], 12, INK, False, False)]], line=1.14); yy += 1.35
    txt(s, MX, yy, lw, 0.3, [[('THE SOLUTION — SWIGGY SARAL', 10, GREEN, True, False)]]); yy += 0.3
    txt(s, MX, yy, lw, 1.1, [[(sl['solution'], 12, INK, False, False)]], line=1.14)
    # decision ask ribbon
    ry = SH-1.15
    rect(s, MX, ry, lw, 0.62, fill=C('FFF4EA'), line=PRIM, lw=1.2, rounded=True, radius=0.08)
    txt(s, MX+0.2, ry+0.09, lw-0.4, 0.5,
        [[('DECISION ASK   ', 10, PRIMD, True, False), (sl['decision_ask'], 11, INK, False, False)]],
        anchor=MSO_ANCHOR.MIDDLE, line=1.05)
    # evidence ledger cards (right)
    ex = MX + lw + 0.45; ew = SW - MX - ex
    txt(s, ex, y+0.1, ew, 0.3, [[('EVIDENCE LEDGER', 10, PRIMD, True, False)]])
    ch = 0.92; cy = y + 0.45
    for it in sl['evidence_ledger']:
        rect(s, ex, cy, ew, ch, fill=SURFACE, line=HAIR, lw=1.0, rounded=True, radius=0.08)
        rect(s, ex, cy+0.12, 0.06, ch-0.24, fill=PRIM)
        txt(s, ex+0.22, cy+0.13, ew-0.4, 0.4, [[(it['value'], 17, INK, True, False, FH)]])
        txt(s, ex+0.22, cy+0.55, ew-0.4, 0.3, [[(it['label'], 10, MUTED, False, False)]])
        cy += ch + 0.14
    footer(s, sl, sl['n']); notes(s, sl)

def s2(s, sl):
    y = header(s, sl)
    n = len(sl['personas']); gap = 0.3
    cw = (SW-2*MX-gap*(n-1))/n; ch = 2.75; cy = y+0.05
    accents = [BLUE, PRIM, RED]
    bgs = [BLUE_BG, C('FCE9DE'), RED_BG]
    for i, p in enumerate(sl['personas']):
        cx = MX + i*(cw+gap)
        rect(s, cx, cy, cw, ch, fill=SURFACE, line=HAIR, lw=1.0, rounded=True, radius=0.05)
        rect(s, cx, cy, cw, 0.5, fill=bgs[i], line=None, rounded=True, radius=0.05)
        rect(s, cx, cy+0.35, cw, 0.16, fill=bgs[i])
        txt(s, cx+0.18, cy+0.1, cw-0.3, 0.35, [[(p['name'], 13, INK, True, False, FH)]])
        txt(s, cx+0.18, cy+0.56, cw-0.3, 0.3, [[(p['tag'], 10, accents[i], True, False)]])
        yy = cy+0.88
        cpl = (cw-0.34)/0.072  # ~chars per line at 9pt
        for lbl, key in [('Context', 'context'), ('Mental model', 'model'), ('Attitude', 'attitude')]:
            txt(s, cx+0.18, yy, cw-0.34, 0.7,
                [[(lbl+': ', 9, accents[i], True, False), (p[key], 9, MUTED, False, False)]], line=1.04)
            import math
            lines = max(1, math.ceil((len(lbl)+2+len(p[key]))/cpl))
            yy += 0.055 + 0.145*lines
        txt(s, cx+0.18, cy+ch-0.32, cw-0.34, 0.3,
            [[('If we fail: ', 9, RED, True, False), (p['risk'], 9, INK, True, False)]])
    # pain points + why
    py = cy + ch + 0.18
    txt(s, MX, py, 6.6, 0.3, [[('PAIN POINTS TODAY', 10, PRIMD, True, False)]])
    bullets(s, MX, py+0.32, 6.6, 1.4, sl['pain_points'], size=10.5, gap=2)
    wx = MX + 6.9
    rect(s, wx, py+0.02, SW-MX-wx, 1.55, fill=C('FFF4EA'), line=PRIM, lw=1.0, rounded=True, radius=0.06)
    txt(s, wx+0.2, py+0.16, SW-MX-wx-0.4, 0.3, [[('WHY IT MATTERS', 10, PRIMD, True, False)]])
    txt(s, wx+0.2, py+0.48, SW-MX-wx-0.4, 1.0, [[(sl['why'], 11.5, INK, False, False)]], line=1.16)
    footer(s, sl, sl['n']); notes(s, sl)

def s_diagram_left(s, sl, right_render):
    y = header(s, sl)
    lw = 5.15
    pic_fit(s, os.path.join(DZ, sl['diagram']), MX, y+0.05, lw, SH-y-0.6, align='center', valign='top')
    rx = MX + lw + 0.4; rw = SW - MX - rx
    right_render(s, sl, rx, y+0.05, rw, SH-y-0.65)
    footer(s, sl, sl['n']); notes(s, sl)

def s3(s, sl):
    y = header(s, sl)
    lw = 4.2
    pic_fit(s, os.path.join(DZ, sl['diagram']), MX, y+0.05, lw, SH-y-1.15, align='center', valign='top')
    rx = MX+lw+0.35; rw = SW-MX-rx
    txt(s, rx, y+0.05, rw, 0.3, [[('WHERE IT FAILS DPDP', 10, PRIMD, True, False)]])
    rows = [['DPDP principle', 'Current state', 'Sev']]
    for g in sl['gaps']:
        rows.append([g['p'], g['s'], g['sev']])
    table(s, rx, y+0.36, rw, rows, [2.0, 3.2, 0.9], fs=9, hfs=9, rh=0.44, hh=0.32,
          cell_cols={2: lambda v: sev_color(v)[0]})
    # risks strip
    ry = SH-0.95
    txt(s, MX, ry-0.28, SW-2*MX, 0.3, [[('FRICTION & RISK POINTS', 10, PRIMD, True, False)]])
    rk = sl['risks']; n=len(rk); gap=0.2; cw=(SW-2*MX-gap*(n-1))/n
    for i, r in enumerate(rk):
        cx = MX+i*(cw+gap)
        rect(s, cx, ry, cw, 0.56, fill=SURFACE, line=HAIR, lw=1.0, rounded=True, radius=0.08)
        txt(s, cx+0.12, ry+0.07, cw-0.2, 0.44, [[(r, 9, INK, False, False)]], line=1.02,
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s, sl, sl['n']); notes(s, sl)

def s4(s, sl):
    y = header(s, sl)
    # right diagram
    dw = 3.4
    dx = SW-MX-dw
    pic_fit(s, os.path.join(DZ, sl['diagram']), dx, y+0.05, dw, SH-y-0.6, align='center', valign='top')
    # left pillars 2-col grid
    lw = dx - MX - 0.35
    txt(s, MX, y+0.02, lw, 0.3, [[('FIVE PILLARS', 10, PRIMD, True, False)]])
    ps = sl['pillars']; cols=2; gap=0.22
    cw=(lw-gap)/cols; ch=0.78; vgap=0.11; gy=y+0.34
    for i,p in enumerate(ps[:4]):
        r=i//cols; c=i%cols
        cx=MX+c*(cw+gap); cy=gy+r*(ch+vgap)
        card(s, cx, cy, cw, ch, p['h'], p['d'], accent=PRIM, head_size=11, body_size=9.5)
    # 5th pillar full width
    cy = gy+2*(ch+vgap)
    card(s, MX, cy, lw, 0.66, ps[4]['h'], ps[4]['d'], accent=GREEN, hcol=INK, head_size=11, body_size=9.5)
    # alternatives
    ay = cy+0.8
    txt(s, MX, ay, lw, 0.3, [[('WHY NOT THE ALTERNATIVES', 10, PRIMD, True, False)]])
    rows=[['Alternative', 'Rejected because']]
    for a in sl['alternatives']:
        rows.append([a['a'], a['r']])
    table(s, MX, ay+0.3, lw, rows, [2.4, 3.0], fs=8.5, hfs=8.5, rh=0.28, hh=0.26)
    footer(s, sl, sl['n']); notes(s, sl)

def s5(s, sl):
    y = header(s, sl)
    # full width kpi band
    bw = SW-2*MX
    pic_fit(s, os.path.join(DZ, sl['diagram']), MX, y+0.03, bw, 1.55, align='center', valign='top')
    cy = y+0.03+1.66
    bottom = SH-0.5
    # north star card
    ns = sl['north_star']
    nsw = 4.3
    rect(s, MX, cy, nsw, bottom-cy, fill=GREEN_BG, line=GREEN, lw=1.2, rounded=True, radius=0.05)
    txt(s, MX+0.2, cy+0.14, nsw-0.4, 0.3, [[('★ NORTH STAR METRIC', 10, GREEN, True, False)]])
    txt(s, MX+0.2, cy+0.44, nsw-0.4, 0.5, [[(ns['name'], 14, INK, True, False, FH)]], line=1.0)
    txt(s, MX+0.2, cy+0.98, nsw-0.4, 0.9, [[(ns['def'], 10, MUTED, False, False)]], line=1.08)
    txt(s, MX+0.2, bottom-0.42, nsw-0.4, 0.4, [[('Target: ', 11, GREEN, True, False), (ns['target'], 11, INK, True, False)]])
    # right column with tracked cursor
    rx = MX+nsw+0.35; rw=SW-MX-rx; c=cy
    txt(s, rx, c, rw, 0.3, [[('SUPPORTING METRICS', 10, PRIMD, True, False)]]); c+=0.3
    rows=[['Metric', 'Baseline', 'Target']]
    for m in sl['supporting']:
        rows.append([m['m'], m['b'], m['t']])
    table(s, rx, c, rw, rows, [2.3,1.4,1.7], fs=9, hfs=9, rh=0.3, hh=0.28)
    c += 0.28 + 0.3*len(sl['supporting']) + 0.14
    txt(s, rx, c, rw, 0.3, [[('GUARDRAILS   ', 10, RED, True, False),
        ('   ·   '.join('%s: %s'%(g['m'],g['t']) for g in sl['guardrails']), 9, MUTED, False, False)]], line=1.05); c+=0.36
    txt(s, rx, c, rw, 0.3, [[('SUCCESS CRITERIA', 10, GREEN, True, False)]]); c+=0.26
    bullets(s, rx, c, rw, bottom-c, sl['success'], size=9, gap=1, mcol=GREEN, marker='✓  ')
    footer(s, sl, sl['n']); notes(s, sl)

def s6(s, sl):
    def right(s, sl, rx, ry, rw, rh):
        txt(s, rx, ry, rw, 0.3, [[('DATA ANALYSIS PLAN', 10, PRIMD, True, False)]])
        bullets(s, rx, ry+0.3, rw, 2.3, sl['analysis_plan'], size=9.8, gap=3)
        y2 = ry+2.55
        rect(s, rx, y2, rw, 0.9, fill=C('FFF4EA'), line=PRIM, lw=1.0, rounded=True, radius=0.05)
        txt(s, rx+0.16, y2+0.1, rw-0.3, 0.3, [[('ROOT-CAUSE METHOD', 9.5, PRIMD, True, False)]])
        txt(s, rx+0.16, y2+0.36, rw-0.3, 0.6, [[(sl['root_cause'], 9.3, INK, False, False)]], line=1.06)
        y3 = y2+1.02
        txt(s, rx, y3, rw, 0.3, [[('NEXT STEPS', 10, GREEN, True, False)]])
        bullets(s, rx, y3+0.28, rw, 1.0, sl['next_steps'], size=9.3, gap=2, mcol=GREEN)
    s_diagram_left(s, sl, right)

def s7(s, sl):
    y = header(s, sl)
    # main KPI + 4 tiles row
    mk = sl['main_kpi']; scol, sbg = status_color(mk['status'])
    mw = 3.2
    rect(s, MX, y+0.05, mw, 1.6, fill=sbg, line=scol, lw=1.2, rounded=True, radius=0.05)
    txt(s, MX+0.18, y+0.16, mw-0.3, 0.3, [[('MAIN KPI · '+mk['metric'], 9.5, scol, True, False)]])
    txt(s, MX+0.18, y+0.45, mw-0.3, 0.6, [[(mk['value'], 34, INK, True, False, FH)]])
    txt(s, MX+0.18, y+1.05, mw-0.3, 0.3, [[('Target '+mk['target']+'   ·   '+mk['trend'], 9.5, scol, True, False)]])
    txt(s, MX+0.18, y+1.32, mw-0.3, 0.3, [[(mk['status'], 10, scol, True, False)]])
    # tiles
    tx = MX+mw+0.3; tw=(SW-MX-tx-0.3*3)/4 + 0
    n=4; gap=0.22; tcw=(SW-MX-tx-gap*(n-1))/n
    for i,t in enumerate(sl['tiles']):
        cx=tx+i*(tcw+gap); tc,tbg=status_color(t['status'])
        rect(s, cx, y+0.05, tcw, 1.6, fill=SURFACE, line=HAIR, lw=1.0, rounded=True, radius=0.06)
        rect(s, cx, y+0.05, tcw, 0.1, fill=tc, rounded=True, radius=0.06)
        txt(s, cx+0.14, y+0.26, tcw-0.24, 0.5, [[(t['m'], 9.3, MUTED, True, False)]], line=1.0)
        txt(s, cx+0.14, y+0.72, tcw-0.24, 0.4, [[(t['v'], 21, INK, True, False, FH)]])
        txt(s, cx+0.14, y+1.18, tcw-0.24, 0.4, [[(t['sub'], 8.5, tc, False, False)]], line=1.0)
    # segmentation table
    sy=y+1.85
    txt(s, MX, sy, 8, 0.3, [[('SEGMENTATION — platform × geography × cohort', 10, PRIMD, True, False)]])
    rows=[['Segment','VCCR','Onboard','Opt-in','Note']]
    for g in sl['segmentation']:
        rows.append([g['seg'],g['vccr'],g['onb'],g['opt'],g['note']])
    table(s, MX, sy+0.3, 8.05, rows, [2.4,0.8,1.0,0.8,3.0], fs=8.7, hfs=8.7, rh=0.34, hh=0.3)
    # alerts
    ax=MX+8.3; aw=SW-MX-ax
    txt(s, ax, sy, aw, 0.3, [[('ALERT RULES', 10, RED, True, False)]])
    bullets(s, ax, sy+0.3, aw, 2.3, sl['alerts'], size=9, gap=4, mcol=RED, marker='▸  ')
    footer(s, sl, sl['n']); notes(s, sl)

PHASES=[("Phase 0 · Wk 1-3","Legal & Design foundation","12-lang notices · DPIA · Ledger backend"),
        ("Phase 1 · Wk 4-7","Controlled A/B","5% new users · Bangalore + Mumbai"),
        ("Phase 2 · Wk 8-10","Regional pilot","25% · 8 cities · Guardian mode"),
        ("Phase 3 · Wk 11-14","Full rollout + backfill","100% new · re-consent campaign")]
GATES=["Gate: legal sign-off · DPIA · pentest",
       "Gate: no P0 7d · onboarding ±4pp · VCCR >85%",
       "Gate: ledger p99 <300ms · guardian >60%",
       "Exit: VCCR >95% · NPS +10 · audit clean"]

def native_timeline(s, x, y, w):
    n=4; gap=0.28; cw=(w-gap*(n-1))/n; ch=1.05
    for i in range(n):
        cx=x+i*(cw+gap)
        rect(s, cx, y, cw, ch, fill=C('FFF4EA'), line=PRIM, lw=1.1, rounded=True, radius=0.06)
        rect(s, cx, y, cw, 0.08, fill=PRIM, rounded=True, radius=0.06)
        txt(s, cx+0.14, y+0.14, cw-0.24, 0.3, [[(PHASES[i][0], 10.5, PRIMD, True, False, FH)]])
        txt(s, cx+0.14, y+0.42, cw-0.24, 0.3, [[(PHASES[i][1], 10, INK, True, False)]])
        txt(s, cx+0.14, y+0.66, cw-0.24, 0.36, [[(PHASES[i][2], 8.7, MUTED, False, False)]], line=1.04)
        # gate caption below
        txt(s, cx+0.02, y+ch+0.05, cw-0.02, 0.4, [[(GATES[i], 8.2, GREEN, False, False)]], line=1.02)
        if i<n-1:
            a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(cx+cw+0.02), Inches(y+ch/2-0.09), Inches(gap-0.04), Inches(0.18))
            a.fill.solid(); a.fill.fore_color.rgb=PRIM; a.line.fill.background(); a.shadow.inherit=False

def s8(s, sl):
    y = header(s, sl)
    bw=SW-2*MX
    native_timeline(s, MX, y+0.08, bw)
    cy=y+0.08+1.6
    ab=sl['ab']
    colw=(SW-2*MX-0.6)/3
    # testable
    card_h = SH-cy-0.55
    rect(s, MX, cy, colw, card_h, fill=GREEN_BG, line=GREEN, lw=1.0, rounded=True, radius=0.05)
    txt(s, MX+0.16, cy+0.12, colw-0.3, 0.3, [[('A/B TESTABLE', 10, GREEN, True, False)]])
    bullets(s, MX+0.16, cy+0.44, colw-0.32, card_h-0.5, ab['testable'], size=9.5, gap=3, mcol=GREEN)
    # non-testable
    nx=MX+colw+0.3
    rect(s, nx, cy, colw, card_h, fill=RED_BG, line=RED, lw=1.0, rounded=True, radius=0.05)
    txt(s, nx+0.16, cy+0.12, colw-0.3, 0.3, [[('NEVER A/B-TESTED', 10, RED, True, False)]])
    bullets(s, nx+0.16, cy+0.44, colw-0.32, 1.3, ab['non_testable'], size=9.5, gap=3, mcol=RED)
    txt(s, nx+0.16, cy+card_h-1.15, colw-0.32, 1.1,
        [[('Compliance shadow test: ', 9, RED, True, False), (ab['handling'], 9, INK, False, False)]], line=1.06)
    # validation
    vx=nx+colw+0.3
    rect(s, vx, cy, colw, card_h, fill=SURFACE, line=HAIR, lw=1.0, rounded=True, radius=0.05)
    txt(s, vx+0.16, cy+0.12, colw-0.3, 0.3, [[('VALIDATION METHODS', 10, PRIMD, True, False)]])
    bullets(s, vx+0.16, cy+0.44, colw-0.32, card_h-0.5, sl['validation'], size=9.5, gap=3, mcol=PRIM)
    footer(s, sl, sl['n']); notes(s, sl)

def s9(s, sl):
    y = header(s, sl)
    lw=8.1
    txt(s, MX, y+0.03, lw, 0.3, [[('RISK REGISTER', 10, PRIMD, True, False)]])
    rows=[['Risk','Impact','Prob','Mitigation → Trade-off']]
    for r in sl['risks']:
        rows.append([r['c'], r['i'], r['p'], r['m']+'  →  '+r['t']])
    table(s, MX, y+0.34, lw, rows, [1.7,0.75,0.75,4.6], fs=8.3, hfs=8.5, rh=0.62, hh=0.3,
          cell_cols={1: lambda v: sev_color(v)[0]})
    # heatmap right
    hx=MX+lw+0.35; hw=SW-MX-hx
    txt(s, hx, y+0.03, hw, 0.3, [[('RISK HEATMAP', 10, RED, True, False)]])
    hm=sl['heatmap']; cy=y+0.36
    cols=['RED_BG','FCE9DE','FEF3C7','SURFACE']
    accs=[RED, PRIM, C('B45309'), MUTED]
    ch=(SH-cy-0.55)/4
    for i,q in enumerate(hm):
        fill=[RED_BG, C('FCE9DE'), C('FEF3C7'), SURFACE][i]
        rect(s, hx, cy+i*(ch+0.1), hw, ch, fill=fill, line=accs[i], lw=1.0, rounded=True, radius=0.06)
        txt(s, hx+0.14, cy+i*(ch+0.1)+0.08, hw-0.26, 0.3, [[(q['q'], 9.3, accs[i], True, False)]])
        txt(s, hx+0.14, cy+i*(ch+0.1)+0.34, hw-0.26, ch-0.4, [[(q['items'], 9.3, INK, False, False)]], line=1.05)
    footer(s, sl, sl['n']); notes(s, sl)

def s10(s, sl):
    y = header(s, sl)
    # three pillars
    ps=sl['pillars']; n=3; gap=0.3; cw=(SW-2*MX-gap*(n-1))/n; ch=1.5; cy=y+0.05
    accs=[GREEN, PRIM, BLUE]; bgs=[GREEN_BG, C('FCE9DE'), BLUE_BG]
    for i,p in enumerate(ps):
        cx=MX+i*(cw+gap)
        rect(s, cx, cy, cw, ch, fill=bgs[i], line=accs[i], lw=1.0, rounded=True, radius=0.05)
        txt(s, cx+0.18, cy+0.14, cw-0.34, 0.35, [[(p['h'], 13, INK, True, False, FH)]])
        txt(s, cx+0.18, cy+0.55, cw-0.34, 0.9, [[(p['d'], 10, MUTED, False, False)]], line=1.12)
    # decision ask ribbon
    ay=cy+ch+0.2
    rect(s, MX, ay, SW-2*MX, 0.78, fill=C('FFF4EA'), line=PRIM, lw=1.2, rounded=True, radius=0.06)
    txt(s, MX+0.2, ay+0.1, SW-2*MX-0.4, 0.6,
        [[('DECISION ASK   ', 10.5, PRIMD, True, False), (sl['decision_ask'], 11, INK, False, False)]], line=1.08)
    # next steps table + closing
    ny=ay+0.98
    txt(s, MX, ny, 6.6, 0.3, [[('IMMEDIATE NEXT STEPS', 10, PRIMD, True, False)]])
    rows=[['Owner','Action','Due']]
    for x in sl['next_steps']:
        rows.append([x['o'], x['a'], x['d']])
    table(s, MX, ny+0.3, 6.7, rows, [1.3,3.6,0.8], fs=8.8, hfs=9, rh=0.3, hh=0.28)
    # closing box
    clx=MX+6.95; clw=SW-MX-clx
    rect(s, clx, ny, clw, SH-ny-0.55, fill=INK, line=None, rounded=True, radius=0.05)
    txt(s, clx+0.22, ny+0.2, clw-0.44, SH-ny-0.9,
        [[(sl['closing'], 12.5, C('FFFFFF'), False, True)]], line=1.22, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, sl, sl['n']); notes(s, sl)

BUILDERS={1:s1,2:s2,3:s3,4:s4,5:s5,6:s6,7:s7,8:s8,9:s9,10:s10}
for sl in DECK['slides']:
    s = slide()
    BUILDERS[sl['n']](s, sl)

out = os.path.join(ROOT, 'Swiggy-DPDP-Final-PM-Case.pptx')
prs.save(out)
print('saved', out, '·', len(prs.slides.__iter__.__self__._sldIdLst), 'slides')
